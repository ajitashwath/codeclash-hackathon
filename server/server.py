from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from google import genai
import json
import uuid
import os
from datetime import datetime
import logging
from dotenv import load_dotenv
from io import BytesIO
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app, origins=["http://localhost:3000", "http://127.0.0.1:3000", "http://localhost:3001"])

# Configure Gemini AI (Google GenAI)
GOOGLE_API_KEY = os.getenv('GOOGLE_API_KEY')
if GOOGLE_API_KEY and GOOGLE_API_KEY != 'your_gemini_api_key_here':
    client = genai.Client()
    logger.info("✅ Gemini AI client configured successfully")
else:
    logger.warning("⚠️ GOOGLE_API_KEY not found or not configured")
    client = None

# In-memory storage for demo (use database in production)
presentations = {}
conversations = {}

# Color themes mapping
COLOR_THEMES = {
    "red": {
        "background": "#fee2e2",
        "primary": "#dc2626",
        "secondary": "#fca5a5",
        "text": "#7f1d1d",
        "accent": "#ef4444"
    },
    "blue": {
        "background": "#dbeafe", 
        "primary": "#2563eb",
        "secondary": "#93c5fd",
        "text": "#1e3a8a",
        "accent": "#3b82f6"
    },
    "green": {
        "background": "#dcfce7",
        "primary": "#16a34a", 
        "secondary": "#86efac",
        "text": "#14532d",
        "accent": "#22c55e"
    },
    "yellow": {
        "background": "#fef3c7",
        "primary": "#d97706",
        "secondary": "#fcd34d", 
        "text": "#92400e",
        "accent": "#f59e0b"
    },
    "purple": {
        "background": "#f3e8ff",
        "primary": "#9333ea",
        "secondary": "#c4b5fd",
        "text": "#581c87", 
        "accent": "#a855f7"
    },
    "pink": {
        "background": "#fce7f3",
        "primary": "#db2777",
        "secondary": "#f9a8d4",
        "text": "#831843",
        "accent": "#ec4899"
    },
    "cyan": {
        "background": "#cffafe",
        "primary": "#0891b2",
        "secondary": "#67e8f9", 
        "text": "#164e63",
        "accent": "#06b6d4"
    },
    "lime": {
        "background": "#ecfccb",
        "primary": "#65a30d",
        "secondary": "#bef264",
        "text": "#365314",
        "accent": "#84cc16"
    }
}

def generate_with_gemini(prompt):
    if not client:
        logger.warning("Gemini client not available, using fallback")
        return generate_fallback_content(prompt)
    try:
        enhanced_prompt = (
            f"Create a professional presentation slide based on this request: \"{prompt}\". "
            "Respond ONLY with a valid JSON object, no explanations or extra text. "
            "The JSON should have these keys: title, content, bullet_points, design_theme, layout_type. "
            "Example: "
            "{ "
            "  \"title\": \"Your Slide Title\", "
            "  \"content\": \"A 2-3 sentence summary.\", "
            "  \"bullet_points\": [\"Point 1\", \"Point 2\", \"Point 3\"], "
            "  \"design_theme\": \"professional\", "
            "  \"layout_type\": \"bullet-list\" "
            "} "
            "Respond ONLY with the JSON object."
        )
        response = client.models.generate_content(
            model='gemini-pro',
            contents=enhanced_prompt
        )
        logger.info(f"Gemini raw response: {getattr(response, 'text', str(response))}")
        if response.text is not None:
            text = response.text.strip()
            if text.startswith('{'):
                try:
                    return json.loads(text)
                except json.JSONDecodeError:
                    logger.error(f"Gemini returned invalid JSON: {text}")
                    return generate_fallback_content(prompt)
            else:
                logger.error(f"Gemini did not return JSON. Raw response: {text}")
                return generate_fallback_content(prompt)
        else:
            logger.error("Gemini API returned no text response.")
            return generate_fallback_content(prompt)
    except Exception as e:
        logger.error(f"Gemini API error: {str(e)}")
        return generate_fallback_content(prompt)

def parse_text_response(text, original_prompt):
    lines = text.strip().split('\n')
    title = None
    for line in lines[:3]:
        if line.strip() and not line.startswith('-') and not line.startswith('•'):
            title = line.strip()
            break
    if not title:
        title = f"About {original_prompt.title()}"
    bullet_points = []
    for line in lines:
        if line.strip().startswith('-') or line.strip().startswith('•'):
            bullet_points.append(line.strip().lstrip('-•').strip())
    if not bullet_points:
        bullet_points = [
            f"Overview of {original_prompt}",
            f"Key facts about {original_prompt}",
            f"Why {original_prompt} matters",
            f"Recent developments in {original_prompt}",
            f"What you can do about {original_prompt}"
        ]
    content = None
    for line in lines:
        if (not line.strip().startswith('-') and not line.strip().startswith('•') and line.strip() != title):
            content = line.strip()
            break
    if not content:
        content = f"This slide provides an overview and key points about {original_prompt}."
    return {
        "title": title,
        "content": content,
        "bullet_points": bullet_points[:5],
        "design_theme": "professional",
        "layout_type": "bullet-list"
    }

def generate_fallback_content(prompt):
    return {
        "title": f"About {prompt.title()}",
        "content": f"This slide provides an overview and key points about {prompt}.",
        "bullet_points": [
            f"Overview of {prompt}",
            f"Key facts about {prompt}",
            f"Why {prompt} matters",
            f"Recent developments in {prompt}",
            f"What you can do about {prompt}"
        ],
        "design_theme": "professional",
        "layout_type": "bullet-list"
    }

def convert_to_slide_elements(ai_response, color_theme="blue"):
    elements = []
    y_position = 80
    theme_colors = COLOR_THEMES.get(color_theme, COLOR_THEMES["blue"])
    elements.append({
        "id": f"title_{uuid.uuid4().hex[:8]}",
        "type": "text",
        "content": ai_response.get("title", "Untitled Slide"),
        "x": 50,
        "y": y_position,
        "width": 700,
        "height": 60,
        "style": {
            "fontSize": "24px",
            "fontWeight": "bold",
            "color": theme_colors["primary"]
        }
    })
    y_position += 100
    if ai_response.get("content"):
        elements.append({
            "id": f"content_{uuid.uuid4().hex[:8]}",
            "type": "text",
            "content": ai_response["content"],
            "x": 50,
            "y": y_position,
            "width": 700,
            "height": 80,
            "style": {
                "fontSize": "16px",
                "color": theme_colors["text"],
                "lineHeight": "1.5"
            }
        })
        y_position += 100
    if ai_response.get("bullet_points"):
        bullet_text = "\n".join([f"• {point}" for point in ai_response["bullet_points"]])
        elements.append({
            "id": f"bullets_{uuid.uuid4().hex[:8]}",
            "type": "text",
            "content": bullet_text,
            "x": 50,
            "y": y_position,
            "width": 700,
            "height": len(ai_response["bullet_points"]) * 30 + 20,
            "style": {
                "fontSize": "14px",
                "color": theme_colors["text"],
                "lineHeight": "1.8"
            }
        })
    return elements

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_enhanced_pptx(slides_data):
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is not available")
    
    prs = Presentation()
    
    # Remove default slide if exists
    if len(prs.slides) > 0:
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        for sld in slides:
            xml_slides.remove(sld)
    
    for slide_data in slides_data:
        # Use title and content layout
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide title
        title = slide_data.get('title', 'Slide')
        if slide.shapes.title:
            slide.shapes.title.text = title
            # Format title
            title_paragraph = slide.shapes.title.text_frame.paragraphs[0]
            title_paragraph.font.size = Pt(32)
            title_paragraph.font.bold = True
            
            # Apply color theme if available
            color_theme = slide_data.get('color_theme', 'blue')
            theme_colors = COLOR_THEMES.get(color_theme, COLOR_THEMES['blue'])
            try:
                rgb = hex_to_rgb(theme_colors['primary'])
                title_paragraph.font.color.rgb = RGBColor(*rgb)
            except:
                pass  # Use default color if conversion fails
        
        # Process slide elements
        elements = slide_data.get('elements', [])
        
        # Find content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                content_placeholder = shape
                break
        
        # Collect all text content
        text_elements = []
        bullet_elements = []
        
        for element in elements:
            if element.get('type') == 'text':
                content = element.get('content', '')
                if '•' in content:
                    # This is bullet content
                    bullet_points = [line.strip('• ').strip() for line in content.split('\n') if line.strip()]
                    bullet_elements.extend(bullet_points)
                else:
                    text_elements.append(content)
            elif element.get('type') == 'bulletList':
                content = element.get('content', '')
                bullet_points = [line.strip('• ').strip() for line in content.split('\n') if line.strip()]
                bullet_elements.extend(bullet_points)
        
        # Add content to placeholder
        text_frame = getattr(content_placeholder, 'text_frame', None)
        if text_frame is not None:
            text_frame.clear()
            # Add regular text first
            if text_elements:
                for i, text in enumerate(text_elements):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = text
                    p.font.size = Pt(16)
                    p.space_after = Pt(12)
            # Add bullet points
            if bullet_elements:
                for i, bullet in enumerate(bullet_elements):
                    if not text_elements and i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0  # First level bullet
                    p.font.size = Pt(14)
                    p.space_after = Pt(6)
        
        # Handle table elements
        table_elements = [el for el in elements if el.get('type') == 'table']
        if table_elements:
            for table_element in table_elements:
                table_data = table_element.get('tableData', {})
                rows = table_data.get('rows', 2)
                cols = table_data.get('cols', 2)
                cells = table_data.get('cells', [])
                
                if cells:
                    # Add table to slide
                    left = Inches(1)
                    top = Inches(3)
                    width = Inches(8)
                    height = Inches(2)
                    
                    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                    
                    # Fill table with data
                    for row_idx, row_data in enumerate(cells[:rows]):
                        for col_idx, cell_data in enumerate(row_data[:cols]):
                            if row_idx < len(table.rows) and col_idx < len(table.columns):
                                cell = table.cell(row_idx, col_idx)
                                cell.text = str(cell_data)
                                # Format cell text
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(12)
    
    return prs

@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "gemini_configured": client is not None,
        "pptx_available": PPTX_AVAILABLE,
        "version": "1.0.0"
    })

@app.route('/api/generate-slide', methods=['POST'])
def generate_slide():
    """Generate a single slide from prompt"""
    try:
        data = request.get_json()
        prompt = data.get('prompt', '')
        color_theme = data.get('color_theme', 'blue')
        
        if not prompt:
            return jsonify({"error": "Prompt is required"}), 400
        
        logger.info(f"Generating slide for: {prompt} with color theme: {color_theme}")
        
        # Generate content with Gemini
        ai_response = generate_with_gemini(prompt)
        
        # Convert to slide format with color theme
        slide = {
            "id": int(datetime.now().timestamp()),
            "title": ai_response.get("title", "Generated Slide"),
            "content": "",
            "notes": "",
            "elements": convert_to_slide_elements(ai_response, color_theme),
            "theme": ai_response.get("design_theme", "professional"),
            "layout": ai_response.get("layout_type", "bullet-list"),
            "color_theme": color_theme,
            "background_color": COLOR_THEMES.get(color_theme, COLOR_THEMES["blue"])["background"],
            "ai_metadata": {
                "original_prompt": prompt,
                "generated_at": datetime.now().isoformat()
            }
        }
        
        return jsonify({
            "slide": slide,
            "ai_response": ai_response,
            "message": "Slide generated successfully"
        })
        
    except Exception as e:
        logger.error(f"Error generating slide: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/color-themes', methods=['GET'])
def get_color_themes():
    """Get available color themes"""
    try:
        return jsonify({
            "themes": COLOR_THEMES,
            "available_colors": list(COLOR_THEMES.keys())
        })
    except Exception as e:
        logger.error(f"Error getting color themes: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/presentai', methods=['POST'])
def create_presentation():
    """Create or update presentation"""
    try:
        data = request.get_json()
        prompt = data.get('prompt', '')
        slides = data.get('slides', [])
        color_theme = data.get('color_theme', 'blue')
        
        presentation_id = str(uuid.uuid4())
        
        # If prompt provided, generate new slide
        if prompt and prompt != "Manual save":
            logger.info(f"Creating presentation with prompt: {prompt}")
            ai_response = generate_with_gemini(prompt)
            
            new_slide = {
                "id": len(slides) + 1,
                "title": ai_response.get("title", "Generated Slide"),
                "content": "",
                "notes": "",
                "elements": convert_to_slide_elements(ai_response, color_theme),
                "theme": ai_response.get("design_theme", "professional"),
                "layout": ai_response.get("layout_type", "bullet-list"),
                "color_theme": color_theme,
                "background_color": COLOR_THEMES.get(color_theme, COLOR_THEMES["blue"])["background"]
            }
            
            slides.append(new_slide)
        
        # Store presentation
        presentations[presentation_id] = {
            "id": presentation_id,
            "prompt": prompt,
            "slides": slides,
            "default_color_theme": color_theme,
            "created_at": datetime.now().isoformat(),
            "updated_at": datetime.now().isoformat()
        }
        
        logger.info(f"Created presentation {presentation_id} with {len(slides)} slides")
        
        return jsonify({
            "presentation_id": presentation_id,
            "slides": slides,
            "message": "Presentation created successfully"
        })
        
    except Exception as e:
        logger.error(f"Error creating presentation: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/export-pptx', methods=['POST'])
def export_pptx():
    """Export slides as a PPTX file with enhanced formatting"""
    if not PPTX_AVAILABLE:
        return jsonify({
            "error": "PPTX export is not available. Please install python-pptx: pip install python-pptx"
        }), 500
    
    try:
        data = request.get_json()
        slides_data = data.get('slides', [])
        
        if not slides_data:
            return jsonify({"error": "No slides provided"}), 400
        
        logger.info(f"Exporting {len(slides_data)} slides to PPTX")
        
        # Create enhanced PPTX
        prs = create_enhanced_pptx(slides_data)
        
        # Save to BytesIO
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        # Generate filename with timestamp
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"slideflow_presentation_{timestamp}.pptx"
        
        return send_file(
            pptx_io,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        logger.error(f"Error exporting PPTX: {str(e)}")
        return jsonify({"error": f"PPTX export failed: {str(e)}"}), 500

# Error handlers
@app.errorhandler(404)
def not_found(error):
    return jsonify({"error": "Endpoint not found"}), 404

@app.errorhandler(500)
def internal_error(error):
    return jsonify({"error": "Internal server error"}), 500

if __name__ == '__main__':
    print("🚀 Starting SlideFlow Backend Server...")
    print("📊 Frontend should connect to: http://localhost:5000")
    print("🔗 API endpoints available at: http://localhost:5000/api/")
    print(f"🤖 Gemini AI: {'✅ Connected' if client else '❌ Not configured'}")
    print(f"📄 PPTX Export: {'✅ Available' if PPTX_AVAILABLE else '❌ Not available (install python-pptx)'}")
    
    app.run(debug=True, host='0.0.0.0', port=5000)